from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ARTIFACTS_DIR = Path(__file__).resolve().parent
METRICS_PATH = ARTIFACTS_DIR / "metrics.json"
OUTPUT_PATH = ARTIFACTS_DIR / "trace_viewer_engineering_spec.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(11, 19, 43)
SLATE = RGBColor(28, 37, 65)
TEAL = RGBColor(91, 192, 190)
AMBER = RGBColor(244, 211, 94)
OFF_WHITE = RGBColor(247, 247, 242)
MUTED = RGBColor(179, 186, 198)
RED = RGBColor(210, 87, 87)


def load_metrics() -> dict:
    return json.loads(METRICS_PATH.read_text(encoding="utf-8"))


def add_background(slide, color: RGBColor = NAVY) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.1), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.15), Inches(11.5), Inches(0.45))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED


def add_footer(slide, text: str = "Mock Zarr recording analysis") -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.0), Inches(12.0), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_bullets(
    slide,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 20,
    color: RGBColor = OFF_WHITE,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.space_before = Pt(0)
        p.line_spacing = 1.1
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def add_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    accent: RGBColor = TEAL,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SLATE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.25)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        Inches(0.12),
        height,
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), width - Inches(0.35), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE

    add_bullets(
        slide,
        lines,
        left=left + Inches(0.22),
        top=top + Inches(0.58),
        width=width - Inches(0.36),
        height=height - Inches(0.72),
        font_size=13,
    )


def add_stat_band(slide, stats: list[tuple[str, str, RGBColor]], *, top: float) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), top, Inches(12.0), Inches(1.1))
    band.fill.solid()
    band.fill.fore_color.rgb = SLATE
    band.line.fill.background()

    card_w = Inches(3.8)
    gap = Inches(0.2)
    left = Inches(0.85)
    for label, value, color in stats:
        value_box = slide.shapes.add_textbox(left, top + Inches(0.12), card_w, Inches(0.4))
        p = value_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = "Aptos Display"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = color

        label_box = slide.shapes.add_textbox(left, top + Inches(0.54), card_w, Inches(0.25))
        p = label_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED
        left += card_w + gap


def add_endpoint_lane(slide, *, left: float, title: str, subtitle: str, bullets: list[str], accent: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.9), Inches(3.75), Inches(4.4))
    card.fill.solid()
    card.fill.fore_color.rgb = SLATE
    card.line.color.rgb = accent

    title_box = slide.shapes.add_textbox(left + Inches(0.22), Inches(2.1), Inches(3.2), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = accent

    sub_box = slide.shapes.add_textbox(left + Inches(0.22), Inches(2.5), Inches(3.2), Inches(0.35))
    p = sub_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED

    add_bullets(
        slide,
        bullets,
        left=left + Inches(0.18),
        top=Inches(2.95),
        width=Inches(3.25),
        height=Inches(2.95),
        font_size=13,
    )


def add_arrow_label(slide, *, left: float, top: float, text: str) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, Inches(1.0), Inches(0.42))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()
    tf = arrow.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = NAVY


def add_arch_box(slide, *, left: float, top: float, width: float, height: float, title: str, body: str, accent: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = SLATE
    box.line.color.rgb = accent

    title_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.1), width - Inches(0.2), Inches(0.25))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = accent

    body_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.38), width - Inches(0.22), height - Inches(0.48))
    p = body_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.color.rgb = OFF_WHITE


def format_bytes(n: int) -> str:
    units = ("B", "KiB", "MiB", "GiB")
    size = float(n)
    unit = units[0]
    for unit in units:
        if size < 1024 or unit == units[-1]:
            break
        size /= 1024
    return f"{size:.2f} {unit}"


def build_cover(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Browser-Based Trace Viewer", "Engineering spec presentation for the SWE coding challenge")

    headline = slide.shapes.add_textbox(Inches(0.68), Inches(1.75), Inches(7.0), Inches(1.2))
    tf = headline.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "The core problem is data delivery, not front-end drawing."
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE

    add_bullets(
        slide,
        [
            "Scientists need full-recording context and high-resolution detail with smooth pan and zoom.",
            "The UI should never blank while the next viewport is loading.",
            "Design choice: overview summaries for movement, raw slices only when pixel density justifies them.",
        ],
        left=Inches(0.7),
        top=Inches(3.0),
        width=Inches(6.4),
        height=Inches(2.2),
        font_size=18,
    )

    dataset = metrics["dataset"]
    add_stat_band(
        slide,
        [
            ("Channels", str(dataset["number_of_channels"]), TEAL),
            ("Sample Rate", f'{int(dataset["sample_rate_hz"])} Hz', AMBER),
            ("Duration", "1.5 h", TEAL),
        ],
        top=Inches(5.7),
    )

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(1.75), Inches(4.2), Inches(4.7))
    callout.fill.solid()
    callout.fill.fore_color.rgb = SLATE
    callout.line.color.rgb = AMBER
    callout.line.width = Pt(1.5)

    title_box = slide.shapes.add_textbox(Inches(8.35), Inches(2.0), Inches(3.6), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Working Thesis"
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AMBER

    add_bullets(
        slide,
        [
            "Treat navigation and inspection as different products.",
            "Optimize `current_data`; treat `voltage_data` as metadata.",
            "Host canonical Zarr in S3 and precompute overview pyramids.",
        ],
        left=Inches(8.3),
        top=Inches(2.55),
        width=Inches(3.5),
        height=Inches(2.8),
        font_size=15,
    )
    add_footer(slide)


def build_scope(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Scope", "What this spec goes deep on")
    add_bullets(
        slide,
        [
            "Use small experiments on the mock recording to decide where raw delivery stops making sense.",
            "Design the service around metadata, overview envelopes, and narrow raw detail reads.",
            "Show a concrete AWS deployment and the operational tradeoffs that matter most.",
            "Keep UI work minimal except where it affects data delivery and perceived latency.",
        ],
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(7.2),
        height=Inches(3.5),
        font_size=22,
    )

    add_card(
        slide,
        left=Inches(8.3),
        top=Inches(1.8),
        width=Inches(3.6),
        height=Inches(1.7),
        title="In Scope",
        lines=[
            "Latency and object fan-out",
            "Caching and preprocessing",
            "Hosting and rollout sequence",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        left=Inches(8.3),
        top=Inches(3.8),
        width=Inches(3.6),
        height=Inches(1.7),
        title="Out Of Scope",
        lines=[
            "Bespoke waveform rendering engine",
            "Deep UI design work",
            "Full auth and tenancy design",
        ],
        accent=AMBER,
    )
    add_footer(slide)


def build_data_findings(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "What The Data Taught Me", "The mock file changed the design")

    current_layout = metrics["layout"]["current_data"]
    voltage_layout = metrics["layout"]["voltage_data"]
    dataset = metrics["dataset"]

    add_card(
        slide,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.8),
        height=Inches(3.7),
        title="current_data",
        lines=[
            f'Chunks {current_layout["chunks"]}, shards {current_layout["shards"]}',
            f'{current_layout["data_object_count"]} data objects',
            f'{format_bytes(current_layout["compressed_bytes"])} compressed',
            "Dominates latency-sensitive interactive reads",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        left=Inches(6.95),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(3.7),
        title="voltage_data",
        lines=[
            f'Chunks {voltage_layout["chunks"]}, shards {voltage_layout["shards"]}',
            f'{voltage_layout["data_object_count"]} data objects',
            f'{format_bytes(voltage_layout["compressed_bytes"])} compressed',
            "Effectively metadata in this mock dataset",
        ],
        accent=AMBER,
    )

    add_stat_band(
        slide,
        [
            ("Total Raw Volume", format_bytes(dataset["total_raw_bytes"]), TEAL),
            ("current_data Objects", str(current_layout["data_object_count"]), AMBER),
            ("voltage_data Objects", str(voltage_layout["data_object_count"]), TEAL),
        ],
        top=Inches(5.95),
    )
    add_footer(slide)


def build_experiments(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Focused Experiments", "A small benchmark set was enough to set the first cutoff")

    window_models = {item["scenario"]: item for item in metrics["derived"]["window_models"]}
    benchmarks = {item["scenario"]: item for item in metrics["benchmarks"]}

    add_card(
        slide,
        left=Inches(0.8),
        top=Inches(1.9),
        width=Inches(3.7),
        height=Inches(3.85),
        title="1 second",
        lines=[
            f'Payload: {format_bytes(window_models["1s"]["raw_payload_bytes_all_channels"])} raw',
            f'Warm read: {benchmarks["1s"]["modes"]["all_channels"]["warm_summary_ms"]["average_ms"]} ms',
            "Raw is acceptable here",
        ],
        accent=TEAL,
    )
    add_card(
        slide,
        left=Inches(4.8),
        top=Inches(1.9),
        width=Inches(3.7),
        height=Inches(3.85),
        title="10 seconds",
        lines=[
            f'Payload: {format_bytes(window_models["10s"]["raw_payload_bytes_all_channels"])} raw',
            f'Samples / px: {window_models["10s"]["samples_per_pixel"]}',
            "Already envelope territory",
        ],
        accent=AMBER,
    )
    add_card(
        slide,
        left=Inches(8.8),
        top=Inches(1.9),
        width=Inches(3.7),
        height=Inches(3.85),
        title="5 minutes",
        lines=[
            f'Payload: {format_bytes(window_models["5min"]["raw_payload_bytes_all_channels"])} raw',
            f'Warm read: {benchmarks["5min"]["modes"]["all_channels"]["warm_summary_ms"]["average_ms"]} ms',
            "Too heavy for responsive viewport refresh",
        ],
        accent=RED,
    )

    full_overview = metrics["derived"]["full_recording_overview"]
    cutoff = metrics["derived"]["raw_detail_cutoff_seconds"]
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.8))
    note.fill.solid()
    note.fill.fore_color.rgb = SLATE
    note.line.fill.background()
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        f'Full overview drops from {format_bytes(full_overview["raw_payload_bytes_all_channels"])} raw '
        f'to {format_bytes(full_overview["envelope_payload_bytes_all_channels"])} as a 1200 px min/max envelope. '
        f'Initial raw cutoff: {cutoff} s per channel.'
    )
    run.font.name = "Aptos"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE
    add_footer(slide)


def build_delivery_model(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Delivery Model", "Two lanes, and never a blank viewport")

    add_endpoint_lane(
        slide,
        left=Inches(0.8),
        title="/metadata",
        subtitle="Shape, scales, duration, default window",
        bullets=[
            "Recording dimensions",
            "Unit scales and offsets",
            "Default channels and window",
            "Overview levels available",
        ],
        accent=TEAL,
    )
    add_endpoint_lane(
        slide,
        left=Inches(4.75),
        title="/overview",
        subtitle="Cacheable min/max envelopes sized to viewport width",
        bullets=[
            "Optimized for navigation",
            "Small stable payloads",
            "Good CDN/cache behavior",
            "Precomputed pyramid source",
        ],
        accent=AMBER,
    )
    add_endpoint_lane(
        slide,
        left=Inches(8.7),
        title="/detail",
        subtitle="Raw `int16` only for narrow windows",
        bullets=[
            "Raw under the cutoff",
            "Envelope otherwise",
            "Binary transport preferred",
            "Keep old frame until replacement is ready",
        ],
        accent=TEAL,
    )

    add_arrow_label(slide, left=Inches(3.9), top=Inches(3.7), text="zoom")
    add_arrow_label(slide, left=Inches(7.85), top=Inches(3.7), text="inspect")
    add_footer(slide)


def build_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "AWS Architecture", "Concrete hosting choices for a first production version")

    add_arch_box(slide, left=Inches(0.8), top=Inches(2.1), width=Inches(2.2), height=Inches(1.2), title="Browser", body="Viewer loads static assets and requests overview/detail windows.", accent=TEAL)
    add_arch_box(slide, left=Inches(3.3), top=Inches(2.1), width=Inches(2.2), height=Inches(1.2), title="CloudFront", body="Serves the app and caches overview responses aggressively.", accent=AMBER)
    add_arch_box(slide, left=Inches(5.8), top=Inches(2.1), width=Inches(2.25), height=Inches(1.2), title="ALB + ECS", body="FastAPI on Fargate handles metadata and narrow detail reads.", accent=TEAL)
    add_arch_box(slide, left=Inches(8.35), top=Inches(1.45), width=Inches(2.2), height=Inches(1.2), title="S3 Zarr", body="Canonical sharded recordings.", accent=TEAL)
    add_arch_box(slide, left=Inches(8.35), top=Inches(3.0), width=Inches(2.2), height=Inches(1.2), title="S3 Pyramids", body="Precomputed overview levels.", accent=AMBER)
    add_arch_box(slide, left=Inches(10.8), top=Inches(2.2), width=Inches(1.7), height=Inches(1.2), title="Batch", body="Builds and refreshes pyramids.", accent=RED)

    add_arrow_label(slide, left=Inches(2.95), top=Inches(2.45), text="cache")
    add_arrow_label(slide, left=Inches(5.45), top=Inches(2.45), text="serve")
    add_arrow_label(slide, left=Inches(7.95), top=Inches(1.8), text="raw")
    add_arrow_label(slide, left=Inches(7.95), top=Inches(3.35), text="env")
    add_arrow_label(slide, left=Inches(10.45), top=Inches(2.55), text="build")

    add_bullets(
        slide,
        [
            "This same pattern ports cleanly to GCP or Azure equivalents.",
            "The key separation is canonical storage versus latency-friendly summarized storage.",
        ],
        left=Inches(0.85),
        top=Inches(5.35),
        width=Inches(8.5),
        height=Inches(1.1),
        font_size=16,
    )
    add_footer(slide)


def build_risks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Hard Parts And Tradeoffs", "These are the engineering risks worth spending time on")

    add_card(slide, left=Inches(0.8), top=Inches(1.8), width=Inches(2.8), height=Inches(2.0), title="Shard Fan-Out", lines=["Cross-shard all-channel reads can double object fetches from 48 to 96.", "Tail latency rises fast once the time window spans multiple shards."], accent=RED)
    add_card(slide, left=Inches(3.95), top=Inches(1.8), width=Inches(2.8), height=Inches(2.0), title="Precompute vs On-Demand", lines=["Too little preprocessing hurts panning latency.", "Too much preprocessing increases build time and storage cost."], accent=AMBER)
    add_card(slide, left=Inches(7.1), top=Inches(1.8), width=Inches(2.8), height=Inches(2.0), title="Admission Control", lines=["Wide detail queries should not starve interactive users.", "Payload and concurrency budgets need hard limits."], accent=TEAL)
    add_card(slide, left=Inches(10.25), top=Inches(1.8), width=Inches(2.25), height=Inches(2.0), title="Cache Correctness", lines=["Cache keys must include recording, channels, window, width, and level."], accent=AMBER)

    add_card(slide, left=Inches(1.5), top=Inches(4.2), width=Inches(4.3), height=Inches(1.7), title="UX Continuity", lines=["Blank redraws during zoom are a product failure even if server latency looks fine.", "Keep the previous frame visible until the next payload is decoded."], accent=TEAL)
    add_card(slide, left=Inches(6.2), top=Inches(4.2), width=Inches(5.0), height=Inches(1.7), title="Observability", lines=["Collect p50/p95/p99 latency, cache hit rates, object-store bytes read, and request rejection rates from day one."], accent=RED)
    add_footer(slide)


def build_plan(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Build Sequence", "Pragmatic rollout with focused follow-up experiments")

    add_card(slide, left=Inches(0.8), top=Inches(2.0), width=Inches(3.75), height=Inches(3.1), title="V1", lines=["Metadata endpoint", "Raw detail path", "One overview envelope path", "Binary transport", "Optimistic client swap with no blank states"], accent=TEAL)
    add_card(slide, left=Inches(4.8), top=Inches(2.0), width=Inches(3.75), height=Inches(3.1), title="V2", lines=["Multiresolution pyramid generation", "CloudFront caching", "Adjacent-window prefetch", "Real object-store performance tests"], accent=AMBER)
    add_card(slide, left=Inches(8.8), top=Inches(2.0), width=Inches(3.75), height=Inches(3.1), title="V3", lines=["SLO tuning", "Admission control", "Auth and tenancy work", "Cost dashboards", "Regression benchmarking"], accent=RED)

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.55), Inches(11.75), Inches(0.95))
    note.fill.solid()
    note.fill.fore_color.rgb = SLATE
    note.line.fill.background()
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Next experiments: S3-backed reads, browser-side binary decode cost, and pyramid build cost per recording."
    run.font.name = "Aptos"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = OFF_WHITE
    add_footer(slide)


def build_recommendation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Recommendation", "The shortest path to a credible first version")

    headline = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(0.7))
    p = headline.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Navigation and inspection are different products."
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = AMBER

    add_bullets(
        slide,
        [
            "The mock data materially shaped the design: `current_data` is the hot path, and 10-second windows are already envelope territory.",
            "Ship a narrow, evidence-backed first version rather than a broad speculative platform.",
            "Success means no blank frames, common interactions under about 150 ms, and predictable cost per recording.",
        ],
        left=Inches(0.85),
        top=Inches(2.7),
        width=Inches(7.6),
        height=Inches(2.6),
        font_size=19,
    )

    add_card(
        slide,
        left=Inches(8.65),
        top=Inches(2.3),
        width=Inches(3.6),
        height=Inches(2.7),
        title="Key Numeric Callout",
        lines=[
            "1.21 GiB raw full overview",
            "225 KiB as 1200 px envelope",
            "5625x reduction",
        ],
        accent=TEAL,
    )
    add_footer(slide, "Generated from artifacts/metrics.json")


def build_presentation(metrics: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build_cover(prs, metrics)
    build_scope(prs)
    build_data_findings(prs, metrics)
    build_experiments(prs, metrics)
    build_delivery_model(prs)
    build_architecture(prs)
    build_risks(prs)
    build_plan(prs)
    build_recommendation(prs)
    return prs


def main() -> None:
    metrics = load_metrics()
    prs = build_presentation(metrics)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
